"""
Attachment Handler for Email Intelligence

Handles processing and storage of email attachments.
"""

from typing import Optional, Dict, Any, List
from pathlib import Path
import os
import hashlib
import logging

from sqlalchemy.orm import Session

from ..models import EmailAttachment, EmailMessage
from ..exceptions import AttachmentError, AttachmentSizeError, AttachmentTypeError
from ..config import get_config

logger = logging.getLogger(__name__)


class AttachmentHandler:
    """
    Handles processing and storage of email attachments.
    
    This class handles:
    - Saving attachments to filesystem
    - Extracting text from various file types
    - Validating attachment size and type
    - Managing attachment storage
    """
    
    def __init__(self, session: Optional[Session] = None):
        """
        Initialize the attachment handler.
        
        Args:
            session: SQLAlchemy session (optional)
        """
        self.session = session
        self.config = get_config()
        
        # Ensure attachment directory exists
        self._ensure_attachment_dir()
    
    def _ensure_attachment_dir(self):
        """Ensure the attachment directory exists"""
        try:
            attachments_dir = self.config.attachments_dir
            attachments_dir.mkdir(parents=True, exist_ok=True)
        except Exception as e:
            logger.error(f"Failed to create attachment directory: {str(e)}")
    
    def process_attachment(self, raw_attachment: Dict[str, Any], email_message: EmailMessage) -> Optional[EmailAttachment]:
        """
        Process a raw attachment and create an EmailAttachment instance.
        
        Args:
            raw_attachment: Dictionary with raw attachment data
            email_message: Parent EmailMessage instance
            
        Returns:
            EmailAttachment instance, or None if processing failed
        """
        try:
            # Validate attachment
            self._validate_attachment(raw_attachment)
            
            # Save attachment to filesystem
            storage_path = self._save_attachment(raw_attachment)
            if not storage_path:
                return None
            
            # Extract text from attachment
            extracted_text = self._extract_text(raw_attachment)
            
            # Create EmailAttachment instance
            attachment = EmailAttachment(
                email_id=email_message.id,
                filename=raw_attachment.get('filename', 'unknown'),
                original_filename=raw_attachment.get('filename', 'unknown'),
                content_type=raw_attachment.get('content_type', 'application/octet-stream'),
                file_size=raw_attachment.get('file_size', 0),
                file_hash=raw_attachment.get('file_hash', self._calculate_hash(raw_attachment.get('payload', b''))),
                storage_path=str(storage_path),
                storage_type='filesystem',
                extracted_text=extracted_text,
                is_inline=self._is_inline_attachment(raw_attachment),
            )
            
            return attachment
            
        except AttachmentSizeError as e:
            logger.warning(f"Attachment too large: {str(e)}")
            return None
        except AttachmentTypeError as e:
            logger.warning(f"Unsupported attachment type: {str(e)}")
            return None
        except Exception as e:
            logger.error(f"Failed to process attachment: {str(e)}")
            return None
    
    def _validate_attachment(self, raw_attachment: Dict[str, Any]):
        """
        Validate an attachment.
        
        Args:
            raw_attachment: Dictionary with raw attachment data
            
        Raises:
            AttachmentSizeError: If attachment exceeds size limit
            AttachmentTypeError: If attachment type is not supported
        """
        # Check size
        file_size = raw_attachment.get('file_size', 0)
        max_size = self.config.processing.max_attachment_size_bytes
        
        if file_size > max_size:
            raise AttachmentSizeError(
                f"Attachment {raw_attachment.get('filename', 'unknown')} exceeds maximum size "
                f"({file_size} > {max_size} bytes)"
            )
        
        # Check type
        content_type = raw_attachment.get('content_type', '').lower()
        supported_types = self.config.processing.supported_attachment_types
        
        if content_type and content_type not in supported_types:
            # Check if it's a common type that we should support
            if not self._is_common_type(content_type):
                raise AttachmentTypeError(
                    f"Attachment type {content_type} is not supported"
                )
    
    def _is_common_type(self, content_type: str) -> bool:
        """
        Check if a content type is commonly supported.
        
        Args:
            content_type: MIME content type
            
        Returns:
            bool: True if the type is commonly supported
        """
        common_types = [
            'application/pdf',
            'application/msword',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'application/vnd.ms-excel',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'application/vnd.ms-powerpoint',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'text/plain',
            'text/html',
            'text/csv',
            'image/jpeg',
            'image/png',
            'image/gif',
            'application/zip',
            'application/x-zip-compressed',
        ]
        
        return content_type.lower() in common_types
    
    def _save_attachment(self, raw_attachment: Dict[str, Any]) -> Optional[Path]:
        """
        Save an attachment to the filesystem.
        
        Args:
            raw_attachment: Dictionary with raw attachment data
            
        Returns:
            Path to the saved file, or None if failed
        """
        try:
            # Get attachment data
            payload = raw_attachment.get('payload')
            if not payload:
                return None
            
            filename = raw_attachment.get('filename', 'unknown')
            content_type = raw_attachment.get('content_type', 'application/octet-stream')
            
            # Generate unique filename
            file_hash = raw_attachment.get('file_hash', self._calculate_hash(payload))
            extension = self._get_extension(content_type, filename)
            unique_filename = f"{file_hash[:16]}{extension}"
            
            # Create storage path
            attachments_dir = self.config.attachments_dir
            storage_path = attachments_dir / unique_filename
            
            # Ensure directory exists
            storage_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Save the file
            with open(storage_path, 'wb') as f:
                f.write(payload)
            
            logger.debug(f"Saved attachment to {storage_path}")
            return storage_path
            
        except Exception as e:
            logger.error(f"Failed to save attachment: {str(e)}")
            return None
    
    def _get_extension(self, content_type: str, filename: str) -> str:
        """
        Get file extension based on content type and filename.
        
        Args:
            content_type: MIME content type
            filename: Original filename
            
        Returns:
            File extension (including dot)
        """
        # Try to get from filename first
        if '.' in filename:
            return Path(filename).suffix.lower()
        
        # Map content types to extensions
        type_to_extension = {
            'application/pdf': '.pdf',
            'application/msword': '.doc',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
            'application/vnd.ms-excel': '.xls',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
            'application/vnd.ms-powerpoint': '.ppt',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
            'text/plain': '.txt',
            'text/html': '.html',
            'text/csv': '.csv',
            'image/jpeg': '.jpg',
            'image/png': '.png',
            'image/gif': '.gif',
            'application/zip': '.zip',
            'application/x-zip-compressed': '.zip',
        }
        
        return type_to_extension.get(content_type.lower(), '.bin')
    
    def _extract_text(self, raw_attachment: Dict[str, Any]) -> Optional[str]:
        """
        Extract text from an attachment.
        
        Args:
            raw_attachment: Dictionary with raw attachment data
            
        Returns:
            Extracted text, or None if extraction failed
        """
        try:
            payload = raw_attachment.get('payload')
            if not payload:
                return None
            
            content_type = raw_attachment.get('content_type', '').lower()
            
            # Try text extraction based on content type
            if content_type.startswith('text/'):
                # Text files - just decode
                try:
                    return payload.decode('utf-8')
                except UnicodeDecodeError:
                    return payload.decode('latin-1')
            
            elif content_type == 'application/pdf':
                return self._extract_pdf_text(payload)
            
            elif content_type in [
                'application/msword',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            ]:
                return self._extract_docx_text(payload)
            
            elif content_type in [
                'application/vnd.ms-excel',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            ]:
                return self._extract_excel_text(payload)
            
            elif content_type in [
                'application/vnd.ms-powerpoint',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            ]:
                return self._extract_powerpoint_text(payload)
            
            elif content_type.startswith('image/'):
                if self.config.processing.perform_ocr:
                    return self._extract_image_text(payload)
                else:
                    return f"[Image: {raw_attachment.get('filename', 'unknown')}]"
            
            else:
                # For unknown types, try to decode as text
                try:
                    return payload.decode('utf-8')
                except UnicodeDecodeError:
                    return f"[Binary file: {raw_attachment.get('filename', 'unknown')}]"
            
        except Exception as e:
            logger.warning(f"Failed to extract text from attachment: {str(e)}")
            return f"[Text extraction failed: {raw_attachment.get('filename', 'unknown')}]"
    
    def _extract_pdf_text(self, payload: bytes) -> Optional[str]:
        """Extract text from PDF"""
        try:
            # Try using pdfminer.six
            try:
                from pdfminer.high_level import extract_text
                return extract_text(payload)
            except ImportError:
                logger.warning("pdfminer.six not available for PDF text extraction")
            except Exception as e:
                logger.warning(f"PDF text extraction failed: {str(e)}")
            
            # Try using PyPDF2 as fallback
            try:
                from PyPDF2 import PdfReader
                import io
                reader = PdfReader(io.BytesIO(payload))
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
            except ImportError:
                logger.warning("PyPDF2 not available for PDF text extraction")
            except Exception as e:
                logger.warning(f"PDF text extraction (PyPDF2) failed: {str(e)}")
            
            return "[PDF content - text extraction requires pdfminer.six or PyPDF2]"
            
        except Exception as e:
            logger.error(f"Failed to extract PDF text: {str(e)}")
            return None
    
    def _extract_docx_text(self, payload: bytes) -> Optional[str]:
        """Extract text from DOCX"""
        try:
            try:
                from docx import Document
                import io
                doc = Document(io.BytesIO(payload))
                text = "\n".join([para.text for para in doc.paragraphs])
                return text
            except ImportError:
                logger.warning("python-docx not available for DOCX text extraction")
            except Exception as e:
                logger.warning(f"DOCX text extraction failed: {str(e)}")
            
            return "[Word document - text extraction requires python-docx]"
            
        except Exception as e:
            logger.error(f"Failed to extract DOCX text: {str(e)}")
            return None
    
    def _extract_excel_text(self, payload: bytes) -> Optional[str]:
        """Extract text from Excel"""
        try:
            try:
                import pandas as pd
                import io
                df = pd.read_excel(io.BytesIO(payload))
                return df.to_string()
            except ImportError:
                logger.warning("pandas not available for Excel text extraction")
            except Exception as e:
                logger.warning(f"Excel text extraction failed: {str(e)}")
            
            try:
                from openpyxl import load_workbook
                import io
                wb = load_workbook(io.BytesIO(payload))
                text = ""
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    text += f"Sheet: {sheet}\n"
                    for row in ws.iter_rows(values_only=True):
                        text += "\t".join(str(cell) if cell else "" for cell in row) + "\n"
                    text += "\n"
                return text
            except ImportError:
                logger.warning("openpyxl not available for Excel text extraction")
            except Exception as e:
                logger.warning(f"Excel text extraction (openpyxl) failed: {str(e)}")
            
            return "[Excel spreadsheet - text extraction requires pandas or openpyxl]"
            
        except Exception as e:
            logger.error(f"Failed to extract Excel text: {str(e)}")
            return None
    
    def _extract_powerpoint_text(self, payload: bytes) -> Optional[str]:
        """Extract text from PowerPoint"""
        try:
            try:
                from pptx import Presentation
                import io
                prs = Presentation(io.BytesIO(payload))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            except ImportError:
                logger.warning("python-pptx not available for PowerPoint text extraction")
            except Exception as e:
                logger.warning(f"PowerPoint text extraction failed: {str(e)}")
            
            return "[PowerPoint presentation - text extraction requires python-pptx]"
            
        except Exception as e:
            logger.error(f"Failed to extract PowerPoint text: {str(e)}")
            return None
    
    def _extract_image_text(self, payload: bytes) -> Optional[str]:
        """Extract text from image using OCR"""
        try:
            if not self.config.processing.perform_ocr:
                return "[Image - OCR disabled]"
            
            try:
                from PIL import Image
                import pytesseract
                import io
                
                # Open the image
                img = Image.open(io.BytesIO(payload))
                
                # Perform OCR
                text = pytesseract.image_to_string(img)
                return text
                
            except ImportError:
                logger.warning("pytesseract or PIL not available for OCR")
            except Exception as e:
                logger.warning(f"OCR failed: {str(e)}")
            
            return "[Image - OCR requires pytesseract and PIL]"
            
        except Exception as e:
            logger.error(f"Failed to extract image text: {str(e)}")
            return None
    
    def _is_inline_attachment(self, raw_attachment: Dict[str, Any]) -> bool:
        """
        Check if an attachment is inline (e.g., embedded image).
        
        Args:
            raw_attachment: Dictionary with raw attachment data
            
        Returns:
            bool: True if the attachment is inline
        """
        content_disposition = raw_attachment.get('content_disposition', '').lower()
        return 'inline' in content_disposition
    
    def _calculate_hash(self, data: bytes) -> str:
        """
        Calculate SHA-256 hash of data.
        
        Args:
            data: Data to hash
            
        Returns:
            Hexadecimal hash string
        """
        return hashlib.sha256(data).hexdigest()
    
    def delete_attachment(self, attachment: EmailAttachment) -> bool:
        """
        Delete an attachment from storage.
        
        Args:
            attachment: EmailAttachment instance
            
        Returns:
            bool: True if deletion successful
        """
        try:
            if attachment.storage_path and attachment.storage_type == 'filesystem':
                storage_path = Path(attachment.storage_path)
                if storage_path.exists():
                    storage_path.unlink()
                    logger.info(f"Deleted attachment file: {storage_path}")
                else:
                    logger.warning(f"Attachment file not found: {storage_path}")
            
            # Remove from database
            if self.session and attachment.id:
                self.session.delete(attachment)
                self.session.commit()
            
            return True
            
        except Exception as e:
            logger.error(f"Failed to delete attachment: {str(e)}")
            if self.session:
                self.session.rollback()
            return False
    
    def get_attachment_content(self, attachment: EmailAttachment) -> Optional[bytes]:
        """
        Get the content of an attachment.
        
        Args:
            attachment: EmailAttachment instance
            
        Returns:
            Attachment content as bytes, or None if failed
        """
        try:
            if attachment.storage_path and attachment.storage_type == 'filesystem':
                storage_path = Path(attachment.storage_path)
                if storage_path.exists():
                    with open(storage_path, 'rb') as f:
                        return f.read()
            
            return None
            
        except Exception as e:
            logger.error(f"Failed to read attachment content: {str(e)}")
            return None
